import os
import tempfile
from typing import List, Dict
from pypdf import PdfReader
from pptx import Presentation

def extract_text_from_txt(file) -> str:
    return file.read().decode("utf-8")

def extract_text_from_pdf(file) -> str:
    reader = PdfReader(file)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text

def extract_text_from_pptx(file) -> str:
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def handle_file_upload(uploaded_files) -> List[Dict]:
    docs = []
    for file in uploaded_files:
        if file.name.endswith(".txt"):
            content = extract_text_from_txt(file)
        elif file.name.endswith(".pdf"):
            content = extract_text_from_pdf(file)
        elif file.name.endswith(".pptx"):
            content = extract_text_from_pptx(file)
        else:
            continue
        docs.append({"filename": file.name, "content": content})
    return docs

def clear_uploaded_files():
    # No persistent storage, so nothing to clear
    pass 